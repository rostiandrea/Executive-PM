"""
Fill assets/steering-committee/TemplateSteerco.pptx with fresh data to
produce an updated Steering Committee deck.

This module is meant to be imported and driven from a session that already
has the Jira extraction in hand (see references/steering-committee-pptx.md
for the full workflow and the data shape expected by each function). It is
deliberately NOT a CLI with a hardcoded data source: Jira data in this repo
is pulled live via the Atlassian connector inside a Claude session, not
from a local file, so the natural call site is a short script written for
that session that imports these helpers and passes in the extracted data.

Golden rule (see references/steering-committee-pptx.md): the Iniziative
tables (slides 9-12) must change as little as possible graphically. The
functions below only ever touch cell text and the RAG cell fill there —
never table style, borders, widths, or row height.

Usage sketch:

    import copy, shutil
    from pptx import Presentation
    from build_steerco_pptx import (
        update_cover, update_summary, update_focus_table,
        update_iniziative_tables, RAG_FILL,
    )

    shutil.copy("assets/steering-committee/TemplateSteerco.pptx", "out.pptx")
    prs = Presentation("out.pptx")

    update_cover(prs, month_year="Settembre 2026")
    update_summary(prs, bullets=[...], kpis=(7, 5, 10), chart_categories=[...], chart_series={...})
    update_focus_table(prs, rows=[...])
    update_iniziative_tables(prs, slides_rows=[rows_slide9, rows_slide10, rows_slide11], suspended_rows=[...])

    prs.save("out.pptx")
"""

from __future__ import annotations

import copy
from typing import Iterable, Sequence

from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

# RAG cell fill colors (hex, no '#'), see assets/steering-committee/brand-tokens.json
RAG_FILL = {
    "green": "00B050",
    "yellow": "FDEA7B",
    "red": "C00000",
}

INIZIATIVE_COLUMNS = [
    "ID", "TITLE", "RAG", "%", "PHASE", "RELEASE", "Ultimo Stato", "Prossimi Passi",
]
FOCUS_COLUMNS = [
    "PRIO", "Codice", "Tavolo", "Descrizione", "Stato", "Last Updates", "Next Step", "Live",
]


def _shape_by_name(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    raise KeyError(f"shape {name!r} not found on slide")


def _set_shape_text(shape, text: str):
    """Replace a shape's text, keeping the formatting of its first run."""
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if not para.runs:
        para.text = text
        return
    para.runs[0].text = text
    for extra in para.runs[1:]:
        extra.text = ""
    for extra_para in tf.paragraphs[1:]:
        for run in extra_para.runs:
            run.text = ""


def update_cover(prs, month_year: str, title_lines: Sequence[str] | None = None):
    """Slide 1: update the '<Month Year>' line (and optionally the title)."""
    slide = prs.slides[0]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text
        if "20" in text and any(ch.isdigit() for ch in text):
            _set_shape_text(shape, month_year)
            return
    raise RuntimeError("cover date placeholder not found — check slide 1 shape names")


def update_summary(
    prs,
    bullets: Sequence[str],
    kpi_new: int,
    kpi_completed: int,
    kpi_releasing: int,
    chart_categories: Sequence[str],
    chart_series: dict[str, Sequence[float | None]],
):
    """Slide 4: bullets, 3 KPI numbers, and the native stacked-bar chart.

    chart_series maps series name -> one value per category (use None for
    a gap, matching the template's sparse stacked layout).
    """
    slide = prs.slides[3]

    bullets_shape = _shape_by_name(slide, "Rectangle 15")
    tf = bullets_shape.text_frame
    for i, line in enumerate(bullets):
        if i < len(tf.paragraphs):
            para = tf.paragraphs[i]
            if para.runs:
                para.runs[0].text = line
            else:
                para.text = line
        else:
            para = tf.add_paragraph()
            para.text = line

    _set_shape_text(_shape_by_name(slide, "Rectangle 14"), str(kpi_new))
    _set_shape_text(_shape_by_name(slide, "Rectangle 12"), str(kpi_completed))
    _set_shape_text(_shape_by_name(slide, "Rectangle 19"), str(kpi_releasing))

    chart_shape = next(s for s in slide.shapes if s.has_chart)
    chart_data = CategoryChartData()
    chart_data.categories = list(chart_categories)
    for name, values in chart_series.items():
        chart_data.add_series(name, values)
    chart_shape.chart.replace_data(chart_data)


def _clear_table_to_row_count(table, n_data_rows: int, header_rows: int = 1):
    """Grow/shrink a table's body to exactly n_data_rows by cloning or
    dropping <a:tr> elements, never touching header rows or table style."""
    tbl = table._tbl
    trs = tbl.findall(qn("a:tr"))
    body_trs = trs[header_rows:]
    current = len(body_trs)

    if current < n_data_rows:
        template_tr = body_trs[-1]
        for _ in range(n_data_rows - current):
            new_tr = copy.deepcopy(template_tr)
            tbl.append(new_tr)
    elif current > n_data_rows:
        for tr in body_trs[n_data_rows:]:
            tbl.remove(tr)


def _fill_row(table, row_idx: int, values: Sequence[str], rag_color: str | None = None, rag_col_idx: int = 2):
    row = table.rows[row_idx]
    for col_idx, value in enumerate(values):
        cell = row.cells[col_idx]
        tf = cell.text_frame
        para = tf.paragraphs[0]
        if not para.runs:
            para.text = str(value)
        else:
            para.runs[0].text = str(value)
            for extra in para.runs[1:]:
                extra.text = ""
        for extra_para in tf.paragraphs[1:]:
            for run in extra_para.runs:
                run.text = ""
    if rag_color is not None:
        rag_cell = row.cells[rag_col_idx]
        rag_cell.fill.solid()
        rag_cell.fill.fore_color.rgb = RGBColor.from_string(rag_color)


def update_focus_table(prs, rows: Sequence[Sequence[str]]):
    """Slide 6: Focus Tavoli SIMA. Each row is the 8 FOCUS_COLUMNS values
    (Stato is plain text here, no color coding on this table)."""
    slide = prs.slides[5]
    table_shape = next(s for s in slide.shapes if s.has_table)
    table = table_shape.table
    _clear_table_to_row_count(table, len(rows))
    for i, row_values in enumerate(rows):
        _fill_row(table, i + 1, row_values, rag_color=None)


def update_iniziative_table(prs, slide_index: int, rows: Sequence[dict]):
    """Fill one Iniziative-schema table (slides 9, 10, 11 or 12).

    Each row dict has keys: id, title, rag ('green'|'yellow'|'red'|None),
    pct, phase, release, ultimo_stato, prossimi_passi.
    """
    slide = prs.slides[slide_index]
    table_shape = next(s for s in slide.shapes if s.has_table)
    table = table_shape.table
    _clear_table_to_row_count(table, len(rows))
    for i, row in enumerate(rows):
        values = [
            row["id"], row["title"], "", str(row["pct"]), row["phase"],
            row.get("release", ""), row.get("ultimo_stato", ""), row.get("prossimi_passi", ""),
        ]
        rag_color = RAG_FILL.get(row["rag"]) if row.get("rag") else None
        _fill_row(table, i + 1, values, rag_color=rag_color, rag_col_idx=2)


def update_iniziative_tables(
    prs,
    slides_rows: Sequence[Sequence[dict]],
    suspended_rows: Sequence[dict] | None = None,
):
    """slides_rows: up to 3 lists of row-dicts for slides 9, 10, 11 (in
    corso, paginated). suspended_rows: row-dicts for slide 12."""
    for offset, rows in enumerate(slides_rows):
        update_iniziative_table(prs, 8 + offset, rows)  # slides[8] == slide 9
    if suspended_rows is not None:
        update_iniziative_table(prs, 11, suspended_rows)  # slides[11] == slide 12
